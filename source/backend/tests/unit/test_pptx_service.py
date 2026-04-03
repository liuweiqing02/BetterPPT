from __future__ import annotations

import base64
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx_service import generate_pptx_from_plan


class PPTXServiceTestCase(unittest.TestCase):
    @staticmethod
    def _write_test_png(path: Path) -> None:
        png_bytes = base64.b64decode(
            'iVBORw0KGgoAAAANSUhEUgAAAAIAAAACCAYAAABytg0kAAAAFklEQVR42mP8z/CfAQgwgImB'
            'QjAAAwMAA/0Bv0QAAAAASUVORK5CYII='
        )
        path.write_bytes(png_bytes)

    def test_generate_editable_pptx(self) -> None:
        slide_plan = [
            {
                'title': 'Cover',
                'bullets': ['Background', 'Goal'],
                'page_function': 'cover',
                'style_tokens_json': {
                    'accent_color': '#1B6EF3',
                    'background_color': '#F7FAFF',
                    'text_color': '#1E2B39',
                    'font_family': 'Arial',
                },
            },
            {
                'title': 'Comparison',
                'bullets': ['Option A', 'Option B', 'Pros', 'Cons'],
                'page_function': 'comparison',
                'style_tokens_json': {
                    'accent_color': '#0F766E',
                    'background_color': '#F0FDFA',
                    'text_color': '#134E4A',
                    'font_family': 'Arial',
                },
            },
        ]

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / 'out.pptx'
            summary = generate_pptx_from_plan(slide_plan=slide_plan, output_path=output_path)

            self.assertTrue(output_path.exists())
            self.assertGreater(output_path.stat().st_size, 0)
            self.assertEqual(summary['page_count'], 2)
            self.assertEqual(len(summary['titles']), 2)

            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 2)
            self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_generate_pptx_with_tables_and_images(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            image_path = Path(tmpdir) / 'hero.png'
            self._write_test_png(image_path)

            slide_plan = [
                {
                    'title': 'Assets',
                    'bullets': ['Keep text editable', 'Render objects separately'],
                    'page_function': 'content',
                    'style_tokens_json': {
                        'accent_color': '#0F766E',
                        'background_color': '#F0FDFA',
                        'text_color': '#134E4A',
                        'font_family': 'Arial',
                    },
                    'tables': [
                        {
                            'title': 'Key Metrics',
                            'headers': ['Metric', 'Value'],
                            'rows': [
                                ['Users', '128'],
                                ['Conversion', '4.2%'],
                            ],
                        }
                    ],
                    'images': [
                        {
                            'caption': 'Hero image placeholder',
                            'alt_text': 'Hero image placeholder',
                            'image_path': str(image_path),
                            'path': str(image_path),
                        }
                    ],
                }
            ]

            output_path = Path(tmpdir) / 'assets.pptx'
            summary = generate_pptx_from_plan(slide_plan=slide_plan, output_path=output_path)

            self.assertTrue(output_path.exists())
            self.assertGreater(output_path.stat().st_size, 0)
            self.assertEqual(summary['page_count'], 1)
            self.assertEqual(len(summary['titles']), 1)
            self.assertEqual(summary['template_edit_stats']['slide_stats'][0]['image_fit_strategy'], 'contain_center')

            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            self.assertTrue(any(getattr(shape, 'has_table', False) for shape in slide.shapes))
            self.assertTrue(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes))

            all_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    all_text.append(shape.text)
            joined_text = '\n'.join(all_text)
            self.assertIn('Hero image placeholder', joined_text)

    def test_generate_pptx_downgrades_complex_table_to_summary_card(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / 'complex_table.pptx'
            slide_plan = [
                {
                    'title': 'Complex Table',
                    'bullets': ['Need readable output'],
                    'page_function': 'content',
                    'style_tokens_json': {
                        'accent_color': '#0F766E',
                        'background_color': '#F0FDFA',
                        'text_color': '#134E4A',
                        'font_family': 'Arial',
                    },
                    'tables': [
                        {
                            'title': 'KPI Matrix',
                            'headers': ['Metric', 'Q1', 'Q2', 'Q3', 'Q4'],
                            'rows': [
                                ['Revenue', 'Very long explanatory text that makes the table hard to read', '12', '13', '14'],
                                ['Growth', '8%', '9%', '10%', '11%'],
                            ],
                        }
                    ],
                }
            ]

            summary = generate_pptx_from_plan(slide_plan=slide_plan, output_path=output_path)

            self.assertTrue(output_path.exists())
            self.assertEqual(summary['template_edit_stats']['slide_stats'][0]['downgraded_table_count'], 1)

            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            self.assertFalse(any(getattr(shape, 'has_table', False) for shape in slide.shapes))

            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    texts.append(shape.text)
            joined = '\n'.join(texts)
            self.assertIn('KPI Matrix', joined)
            self.assertIn('Complex table downgraded', joined)

    def test_generate_pptx_from_template_page(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / 'template.pptx'
            output_path = Path(tmpdir) / 'from_template.pptx'

            template_prs = Presentation()
            template_slide = template_prs.slides.add_slide(template_prs.slide_layouts[1])
            template_slide.shapes.title.text = 'Template Title'
            body = template_slide.placeholders[1].text_frame
            body.clear()
            body.paragraphs[0].text = 'Template Body'
            template_prs.save(str(template_path))

            slide_plan = [
                {
                    'title': 'Generated Title',
                    'bullets': ['Line 1', 'Line 2'],
                    'page_function': 'content',
                    'template_page_no': 1,
                    'style_tokens_json': {'font_family': 'Arial'},
                }
            ]

            summary = generate_pptx_from_plan(
                slide_plan=slide_plan,
                output_path=output_path,
                template_path=template_path,
            )

            self.assertTrue(output_path.exists())
            self.assertTrue(summary['template_mode'])
            self.assertGreaterEqual(summary['template_edit_stats']['text_placeholder_edits'], 1)

            rendered = Presentation(str(output_path))
            self.assertEqual(len(rendered.slides), 1)
            texts = []
            for shape in rendered.slides[0].shapes:
                if getattr(shape, 'has_text_frame', False):
                    texts.append(shape.text)
            joined = '\n'.join(texts)
            self.assertIn('Generated Title', joined)
            self.assertIn('Line 1', joined)

    def test_generate_pptx_template_table_shape_is_edited(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / 'template_table.pptx'
            output_path = Path(tmpdir) / 'from_template_table.pptx'

            template_prs = Presentation()
            slide = template_prs.slides.add_slide(template_prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(0, 0, 4000000, 500000)
            title_box.text = 'Template'
            table_shape = slide.shapes.add_table(3, 2, 0, 700000, 5000000, 2000000)
            table_shape.table.cell(0, 0).text = 'H1'
            table_shape.table.cell(0, 1).text = 'H2'
            template_prs.save(str(template_path))

            slide_plan = [
                {
                    'title': 'Table Slide',
                    'bullets': ['bullet'],
                    'template_page_no': 1,
                    'tables': [
                        {
                            'headers': ['Year', 'Revenue'],
                            'rows': [['2024', '10'], ['2025', '12']],
                        }
                    ],
                }
            ]

            summary = generate_pptx_from_plan(slide_plan=slide_plan, output_path=output_path, template_path=template_path)
            self.assertTrue(summary['template_mode'])
            self.assertGreaterEqual(summary['template_edit_stats']['table_placeholder_edits'], 1)

            rendered = Presentation(str(output_path))
            self.assertEqual(len(rendered.slides), 1)
            table_shapes = [shape for shape in rendered.slides[0].shapes if getattr(shape, 'has_table', False)]
            self.assertTrue(table_shapes)
            table = table_shapes[0].table
            self.assertEqual(table.cell(0, 0).text, 'Year')
            self.assertEqual(table.cell(1, 0).text, '2024')

    def test_generate_pptx_template_named_slot_matching_stats(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            template_path = Path(tmpdir) / 'template_named_slots.pptx'
            output_path = Path(tmpdir) / 'from_template_named_slots.pptx'

            template_prs = Presentation()
            slide = template_prs.slides.add_slide(template_prs.slide_layouts[6])
            title_shape = slide.shapes.add_textbox(0, 0, 5000000, 600000)
            title_shape.name = 'slot_title_main'
            title_shape.text = 'Template Title'
            body_shape = slide.shapes.add_textbox(0, 800000, 7000000, 2600000)
            body_shape.name = 'slot_content_main'
            body_shape.text = 'Template Body'
            template_prs.save(str(template_path))

            slide_plan = [
                {
                    'title': 'Named Slot Title',
                    'bullets': ['Point A', 'Point B'],
                    'template_page_no': 1,
                    'layout_schema_json': {
                        'slots': [
                            {'slot_key': 'slot_title_main', 'slot_type': 'text', 'slot_role': 'title'},
                            {'slot_key': 'slot_content_main', 'slot_type': 'text', 'slot_role': 'bullet'},
                        ]
                    },
                }
            ]

            summary = generate_pptx_from_plan(slide_plan=slide_plan, output_path=output_path, template_path=template_path)
            self.assertTrue(summary['template_mode'])
            self.assertGreaterEqual(summary['template_edit_stats']['named_slot_matches'], 2)
            self.assertEqual(len(summary['template_edit_stats']['slide_stats']), 1)
            self.assertGreaterEqual(summary['template_edit_stats']['slide_stats'][0]['named_matches'], 2)


if __name__ == '__main__':
    unittest.main()
