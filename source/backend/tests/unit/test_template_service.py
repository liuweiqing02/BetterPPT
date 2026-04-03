from __future__ import annotations

import atexit
import json
import os
import sys
from types import SimpleNamespace
import tempfile
import unittest
from pathlib import Path
from types import ModuleType
from unittest.mock import Mock, patch

from sqlalchemy import create_engine, select
from sqlalchemy.orm import Session, sessionmaker

from app.core.config import get_settings
from app.db.base import Base
from app.models import File, TemplatePageSchema, TemplateProfile, User

_REPO_ROOT = Path(__file__).resolve().parents[4]
_SERVICES_DIR = _REPO_ROOT / 'source' / 'backend' / 'app' / 'services'
if 'app.services' not in sys.modules:
    services_pkg = ModuleType('app.services')
    services_pkg.__path__ = [str(_SERVICES_DIR)]  # type: ignore[attr-defined]
    sys.modules['app.services'] = services_pkg

from app.services import template_service as template_service_module
from app.services.template_service import analyze_and_persist_template, analyze_template

_GENERATED_REF_DIR: tempfile.TemporaryDirectory[str] | None = None


def _find_reference_pptx() -> Path:
    root = Path(__file__).resolve().parents[4]
    candidates = sorted((root / 'ref').glob('*.pptx'))
    if candidates:
        return candidates[0]

    # CI/public repos may not include private sample files under ref/.
    # Build a tiny deterministic pptx fixture on the fly as a fallback.
    global _GENERATED_REF_DIR
    if _GENERATED_REF_DIR is None:
        _GENERATED_REF_DIR = tempfile.TemporaryDirectory(prefix='betterppt-unit-ref-')
        atexit.register(_GENERATED_REF_DIR.cleanup)
        generated_path = Path(_GENERATED_REF_DIR.name) / 'generated_reference.pptx'
        from pptx import Presentation

        prs = Presentation()
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = 'BetterPPT Unit Template'
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = 'Auto-generated fallback template'

        content = prs.slides.add_slide(prs.slide_layouts[1])
        content.shapes.title.text = 'Content Layout'
        if len(content.placeholders) > 1:
            content.placeholders[1].text = 'bullet one\nbullet two'

        prs.save(generated_path)
    return Path(_GENERATED_REF_DIR.name) / 'generated_reference.pptx'


class TemplateServiceTestCase(unittest.TestCase):
    def _make_session(self) -> Session:
        tmpdir = tempfile.TemporaryDirectory()
        self.addCleanup(tmpdir.cleanup)

        db_path = Path(tmpdir.name) / 'template_service.sqlite3'
        engine = create_engine(
            f'sqlite+pysqlite:///{db_path.as_posix()}',
            connect_args={'check_same_thread': False},
        )
        Base.metadata.create_all(engine)
        self.addCleanup(engine.dispose)

        session_factory = sessionmaker(bind=engine, autocommit=False, autoflush=False, expire_on_commit=False)
        return session_factory()

    def _create_reference_file(self, db: Session) -> File:
        reference_path = _find_reference_pptx()
        user = User(username='template-service-user', email='template-service@example.com', password_hash='x', status=1)
        db.add(user)
        db.flush()

        reference_file = File(
            user_id=user.id,
            file_role='ppt_reference',
            storage_provider='local',
            storage_path=str(reference_path.resolve()),
            filename=reference_path.name,
            ext='pptx',
            mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            file_size=reference_path.stat().st_size,
            checksum_sha256=None,
            status='uploaded',
        )
        db.add(reference_file)
        db.commit()
        db.refresh(reference_file)
        return reference_file

    def test_analyze_template_returns_pages_with_complete_fields_for_real_pptx(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': ''})

            def _fake_embeddings(pages, detail_level, reference_file):
                return (
                    [[0.1, 0.2, 0.3] for _ in pages],
                    {'embedding_mode': 'rule_features', 'embedding_source': 'rule', 'fallback_reason': None},
                )

            def _fake_cluster(vectors, target_cluster_count):
                labels = [1 if index < max(1, len(vectors) // 2) else 2 for index, _ in enumerate(vectors)]
                return labels, 'hierarchical', None

            with patch.object(
                template_service_module,
                '_build_page_embeddings',
                side_effect=_fake_embeddings,
            ), patch.object(
                template_service_module,
                '_cluster_pages_hierarchically',
                side_effect=_fake_cluster,
            ), patch(
                'app.services.template_service.get_settings',
                return_value=settings,
            ):
                result = analyze_template(reference_file, detail_level='balanced', task_no='unit-template-analyze')
            pages = result['pages']

            self.assertGreater(len(pages), 0)
            self.assertEqual(result['detail_level'], 'balanced')
            self.assertEqual(result['reference_file_id'], reference_file.id)
            self.assertEqual(result['analysis_source'], 'pptx_xml')
            self.assertIn(result['parse_status'], {'ok', 'partial_fallback'})

            for page in pages:
                self.assertIn('page_no', page)
                self.assertIn('cluster_label', page)
                self.assertIn('page_function', page)
                self.assertIn('layout_schema_json', page)
                self.assertIn('style_tokens_json', page)
                self.assertIsInstance(page['page_function'], str)
                self.assertTrue(page['page_function'])
                self.assertIsInstance(page['layout_schema_json'], dict)
                self.assertIsInstance(page['style_tokens_json'], dict)
                self.assertTrue(page['layout_schema_json'])
                self.assertTrue(page['style_tokens_json'])

    def test_analyze_and_persist_template_writes_profile_and_page_schemas(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': ''})

            def _fake_embeddings(pages, detail_level, reference_file):
                return (
                    [[0.1, 0.2, 0.3] for _ in pages],
                    {'embedding_mode': 'rule_features', 'embedding_source': 'rule', 'fallback_reason': None},
                )

            def _fake_cluster(vectors, target_cluster_count):
                labels = [1 if index < max(1, len(vectors) // 2) else 2 for index, _ in enumerate(vectors)]
                return labels, 'hierarchical', None

            with patch.object(
                template_service_module,
                '_build_page_embeddings',
                side_effect=_fake_embeddings,
            ), patch.object(
                template_service_module,
                '_cluster_pages_hierarchically',
                side_effect=_fake_cluster,
            ), patch(
                'app.services.template_service.get_settings',
                return_value=settings,
            ):
                result = analyze_and_persist_template(
                    db,
                    reference_file=reference_file,
                    detail_level='detailed',
                    task_no='unit-template-persist',
                )

            profile = db.scalar(select(TemplateProfile).where(TemplateProfile.id == result['profile_id']))
            self.assertIsNotNone(profile)
            assert profile is not None

            pages = list(
                db.scalars(
                    select(TemplatePageSchema)
                    .where(TemplatePageSchema.template_profile_id == profile.id)
                    .order_by(TemplatePageSchema.page_no.asc())
                ).all()
            )

            self.assertEqual(result['page_schemas_count'], len(pages))
            self.assertEqual(profile.total_pages, result['total_pages'])
            self.assertEqual(profile.cluster_count, result['cluster_count'])
            self.assertEqual(profile.summary_json['task_no'], 'unit-template-persist')
            self.assertIsInstance(profile.summary_json.get('page_function_counts'), dict)
            self.assertEqual(len(pages), profile.total_pages)
            self.assertGreater(len(pages), 0)
            self.assertEqual(result['analysis_source'], 'pptx_xml')
            self.assertEqual(result.get('__persisted__'), True)

            for page in pages:
                self.assertEqual(page.template_profile_id, profile.id)
                self.assertIsInstance(page.layout_schema_json, dict)
                self.assertIsInstance(page.style_tokens_json, dict)
                self.assertTrue(page.page_function)
                self.assertTrue(page.layout_schema_json)
                self.assertTrue(page.style_tokens_json)

    def test_analyze_template_skips_llm_enhancement_when_key_is_missing(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': ''})
            mock_call = Mock(side_effect=AssertionError('llm client should not be called when key is missing'))
            fake_llm_module = ModuleType('app.services.llm_service')
            fake_llm_module.call_chat_completions = mock_call  # type: ignore[attr-defined]

            def _fake_embeddings(pages, detail_level, reference_file):
                return (
                    [[0.1, 0.2, 0.3] for _ in pages],
                    {'embedding_mode': 'rule_features', 'embedding_source': 'rule', 'fallback_reason': None},
                )

            def _fake_cluster(vectors, target_cluster_count):
                labels = [1 if index < max(1, len(vectors) // 2) else 2 for index, _ in enumerate(vectors)]
                return labels, 'hierarchical', None

            with patch.dict(sys.modules, {'app.services.llm_service': fake_llm_module}), patch(
                'app.services.template_service.get_settings', return_value=settings
            ), patch.object(
                template_service_module,
                '_build_page_embeddings',
                side_effect=_fake_embeddings,
            ), patch.object(
                template_service_module,
                '_cluster_pages_hierarchically',
                side_effect=_fake_cluster,
            ):
                result = analyze_template(reference_file, detail_level='balanced', task_no='unit-template-llm-off')

            self.assertFalse(result['llm_enhanced'])
            self.assertEqual(result['llm_model'], settings.llm_model)
            self.assertFalse(result['summary_json']['llm_enhanced'])
            self.assertIsNone(result['summary_json']['llm_usage'])
            self.assertIsNone(result['summary_json']['llm_error'])
            self.assertEqual(result['llm_batches_total'], 0)
            self.assertEqual(result['llm_batches_succeeded'], 0)
            mock_call.assert_not_called()

    def test_analyze_template_applies_llm_page_suggestions_when_mocked(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': 'unit-test-key'})
            llm_payload = {
                'pages': [
                    {
                        'page_no': 1,
                        'page_function': 'cover',
                        'layout_suggestions': {
                            'density_hint': 'compact',
                            'title_style': 'hero',
                        },
                        'style_suggestions': {
                            'accent_strategy': 'brand_strip',
                        },
                        'reason': 'front page should stay visually light',
                    }
                ]
            }
            mock_call = Mock(
                return_value=SimpleNamespace(content=json.dumps(llm_payload), usage={'prompt_tokens': 10, 'completion_tokens': 5})
            )
            fake_llm_module = ModuleType('app.services.llm_service')
            fake_llm_module.call_chat_completions = mock_call  # type: ignore[attr-defined]

            def _fake_embeddings(pages, detail_level, reference_file):
                return (
                    [[0.1, 0.2, 0.3] for _ in pages],
                    {'embedding_mode': 'rule_features', 'embedding_source': 'rule', 'fallback_reason': None},
                )

            def _fake_cluster(vectors, target_cluster_count):
                labels = [1 if index < max(1, len(vectors) // 2) else 2 for index, _ in enumerate(vectors)]
                return labels, 'hierarchical', None

            with patch.dict(sys.modules, {'app.services.llm_service': fake_llm_module}), patch(
                'app.services.template_service.get_settings', return_value=settings
            ), patch.object(
                template_service_module,
                '_build_page_embeddings',
                side_effect=_fake_embeddings,
            ), patch.object(
                template_service_module,
                '_cluster_pages_hierarchically',
                side_effect=_fake_cluster,
            ):
                result = analyze_template(reference_file, detail_level='balanced', task_no='unit-template-llm-on')

            self.assertTrue(result['llm_enhanced'])
            self.assertIsInstance(result['llm_usage'], dict)
            self.assertGreaterEqual(result['llm_usage'].get('prompt_tokens', 0), 10)
            self.assertEqual(result['summary_json']['llm_page_suggestions'][0]['page_no'], 1)
            self.assertGreaterEqual(result['llm_batches_total'], 1)
            self.assertGreaterEqual(result['llm_batches_succeeded'], 1)
            self.assertGreater(len(result['pages']), 0)
            first_page = result['pages'][0]
            self.assertEqual(first_page['layout_schema_json']['layout_rules']['density_hint'], 'compact')
            self.assertEqual(first_page['layout_schema_json']['layout_rules']['title_style'], 'hero')
            self.assertEqual(first_page['style_tokens_json']['accent_strategy'], 'brand_strip')
            self.assertGreaterEqual(mock_call.call_count, 1)

    def test_analyze_template_records_embedding_and_clustering_modes_when_optional_dependencies_work(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': ''})

            def _fake_embeddings(pages, detail_level, reference_file):
                vectors = [[float(index + 1), 0.5, 0.25] for index, _ in enumerate(pages)]
                return vectors, {
                    'embedding_mode': 'vision_model',
                    'embedding_source': 'vision+rule',
                    'fallback_reason': None,
                }

            def _fake_cluster(vectors, target_cluster_count):
                labels = [1 if index < max(1, len(vectors) // 2) else 2 for index, _ in enumerate(vectors)]
                return labels, 'hierarchical', None

            with patch.object(template_service_module, '_build_page_embeddings', side_effect=_fake_embeddings), patch.object(
                template_service_module, '_cluster_pages_hierarchically', side_effect=_fake_cluster
            ), patch('app.services.template_service.get_settings', return_value=settings):
                result = analyze_template(reference_file, detail_level='balanced', task_no='unit-template-modes')

            self.assertEqual(result['embedding_mode'], 'vision_model')
            self.assertEqual(result['embedding_source'], 'vision+rule')
            self.assertEqual(result['clustering_mode'], 'hierarchical')
            self.assertIsNone(result['fallback_reason'])
            self.assertEqual(result['summary_json']['embedding_mode'], 'vision_model')
            self.assertEqual(result['summary_json']['clustering_mode'], 'hierarchical')
            self.assertIsNone(result['summary_json']['fallback_reason'])
            self.assertGreaterEqual(result['cluster_count'], 1)
            self.assertGreater(len(result['pages']), 0)

    def test_analyze_template_records_signature_fallback_when_optional_dependencies_are_missing(self) -> None:
        with self._make_session() as db:
            reference_file = self._create_reference_file(db)
            settings = get_settings().model_copy(update={'llm_api_key': ''})

            with patch.object(
                template_service_module,
                '_build_page_embeddings',
                return_value=(
                    [[0.1, 0.2, 0.3] for _ in range(2)],
                    {
                        'embedding_mode': 'rule_features',
                        'embedding_source': 'rule',
                        'fallback_reason': 'vision embedding failed: unavailable',
                    },
                ),
            ), patch.object(
                template_service_module,
                '_cluster_pages_hierarchically',
                return_value=([], None, 'hierarchical clustering unavailable: scipy missing'),
            ), patch('app.services.template_service.get_settings', return_value=settings):
                result = analyze_template(reference_file, detail_level='balanced', task_no='unit-template-fallback')

            self.assertEqual(result['embedding_mode'], 'rule_features')
            self.assertEqual(result['embedding_source'], 'rule')
            self.assertEqual(result['clustering_mode'], 'signature')
            self.assertIn('vision embedding failed', result['fallback_reason'])
            self.assertIn('hierarchical clustering unavailable', result['fallback_reason'])
            self.assertEqual(result['summary_json']['embedding_mode'], 'rule_features')
            self.assertEqual(result['summary_json']['clustering_mode'], 'signature')
            self.assertIn('vision embedding failed', result['summary_json']['fallback_reason'])

    def test_resolve_vision_model_ref_prefers_local_model_path_env(self) -> None:
        with tempfile.TemporaryDirectory(prefix='vision-local-model-') as tempdir:
            local_model = Path(tempdir).resolve()
            mocked_settings = SimpleNamespace(
                template_vision_model='google/vit-base-patch16-224-in21k',
                template_vision_model_path='',
                template_vision_cache_dir='',
            )
            with patch.object(template_service_module, 'get_settings', return_value=mocked_settings), patch.dict(
                os.environ,
                {
                    'BETTERPPT_TEMPLATE_VISION_MODEL_PATH': str(local_model),
                    'BETTERPPT_TEMPLATE_VISION_MODEL': 'google/vit-base-patch16-224-in21k',
                },
            ):
                resolved = template_service_module._resolve_vision_model_ref()
        self.assertEqual(resolved, str(local_model))

    def test_resolve_vision_model_ref_falls_back_to_model_id(self) -> None:
        mocked_settings = SimpleNamespace(
            template_vision_model='google/vit-base-patch16-224-in21k',
            template_vision_model_path='',
            template_vision_cache_dir='',
        )
        with patch.dict(
            os.environ,
            {
                'BETTERPPT_TEMPLATE_VISION_MODEL_PATH': '',
                'BETTERPPT_TEMPLATE_VISION_MODEL': 'google/vit-base-patch16-224-in21k',
            },
            clear=False,
        ), patch.object(template_service_module, '_DEFAULT_VISION_LOCAL_DIR', Path('Z:/_not_exist_model_dir_')), patch.object(
            template_service_module, 'get_settings', return_value=mocked_settings
        ):
            resolved = template_service_module._resolve_vision_model_ref()
        self.assertEqual(resolved, 'google/vit-base-patch16-224-in21k')


if __name__ == '__main__':
    unittest.main()
