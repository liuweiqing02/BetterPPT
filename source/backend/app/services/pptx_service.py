from __future__ import annotations

from collections.abc import Iterable
from copy import deepcopy
from pathlib import Path
import re
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency guard
    raise RuntimeError(
        'python-pptx is required for pptx generation. Install the `python-pptx` package first.'
    ) from exc


class PPTXGenerationError(RuntimeError):
    pass


def _safe_hex_color(value: Any, default_rgb: tuple[int, int, int]) -> RGBColor:
    if isinstance(value, str):
        raw = value.strip().lstrip('#')
        if len(raw) == 6:
            try:
                return RGBColor.from_string(raw.upper())
            except Exception:
                pass
    return RGBColor(*default_rgb)


def _style_tokens(slide: dict[str, Any]) -> dict[str, Any]:
    style = slide.get('style_tokens_json')
    if isinstance(style, dict):
        return style
    return {}


def _coerce_text(value: Any, default: str = '') -> str:
    if value is None:
        return default
    if isinstance(value, str):
        text = value.strip()
        return text or default
    return str(value).strip() or default


def _coerce_items(value: Any) -> list[Any]:
    if not value:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, tuple):
        return list(value)
    if isinstance(value, dict):
        return [value]
    if isinstance(value, Iterable) and not isinstance(value, (str, bytes)):
        return list(value)
    return [value]


def _normalize_bullets(slide: dict[str, Any], slide_index: int) -> list[str]:
    raw_candidates = (
        slide.get('bullets'),
        slide.get('key_points'),
        slide.get('items'),
        slide.get('points'),
        slide.get('outline'),
    )

    bullets: list[str] = []
    for candidate in raw_candidates:
        for item in _coerce_items(candidate):
            text = _coerce_text(item)
            if text:
                bullets.append(text)

    if bullets:
        return bullets

    title = _coerce_text(slide.get('title') or slide.get('heading') or slide.get('page_title'), default=f'Slide {slide_index}')
    summary = _coerce_text(
        slide.get('summary')
        or slide.get('description')
        or slide.get('content')
        or slide.get('body')
        or slide.get('text'),
        default='',
    )
    return [
        f'Key idea: {title}',
        summary or 'Auto-generated placeholder bullet for missing slide_plan bullets.',
    ]


def _normalize_title(slide: dict[str, Any], slide_index: int, title_prefix: str) -> str:
    base_title = _coerce_text(
        slide.get('title')
        or slide.get('heading')
        or slide.get('page_title')
        or slide.get('section_title')
        or slide.get('name'),
        default=f'Slide {slide_index}',
    )
    if title_prefix:
        return f'{title_prefix}{base_title}'
    return base_title


def _normalize_tables(slide: dict[str, Any]) -> list[dict[str, Any]]:
    raw_candidates = (
        slide.get('tables'),
        slide.get('table'),
        slide.get('table_data'),
        slide.get('table_slots'),
    )

    tables: list[dict[str, Any]] = []
    for candidate in raw_candidates:
        for item in _coerce_items(candidate):
            if not isinstance(item, dict):
                continue

            title = _coerce_text(item.get('title') or item.get('name') or item.get('caption') or item.get('slot_key'))
            headers = [_coerce_text(value) for value in _coerce_items(item.get('headers') or item.get('columns')) if _coerce_text(value)]
            rows: list[list[str]] = []
            raw_rows = item.get('rows') or item.get('data') or item.get('items')

            for row in _coerce_items(raw_rows):
                if isinstance(row, dict):
                    row_values = [_coerce_text(value, default='-') for value in row.values()]
                    if row_values:
                        rows.append(row_values)
                    if not headers:
                        headers = [_coerce_text(key) for key in row.keys()]
                elif isinstance(row, (list, tuple)):
                    row_values = [_coerce_text(value, default='-') for value in row]
                    if row_values:
                        rows.append(row_values)
                else:
                    row_text = _coerce_text(row, default='-')
                    if row_text:
                        rows.append([row_text])

            if not rows:
                hint = item.get('hint') or item.get('summary') or item.get('content') or item.get('text')
                hint_text = _coerce_text(hint, default='')
                if hint_text:
                    rows = [[hint_text]]

            if not rows and not headers:
                continue

            if not headers:
                headers = ['Item', 'Value']

            tables.append(
                {
                    'slot_key': _coerce_text(item.get('slot_key') or item.get('key') or title),
                    'title': title,
                    'headers': headers,
                    'rows': rows,
                    'source': _coerce_text(item.get('source') or item.get('content_source')),
                }
            )

    return tables


def _normalize_images(slide: dict[str, Any]) -> list[dict[str, Any]]:
    raw_candidates = (
        slide.get('images'),
        slide.get('image'),
        slide.get('image_slots'),
        slide.get('visuals'),
    )

    images: list[dict[str, Any]] = []
    for candidate in raw_candidates:
        for item in _coerce_items(candidate):
            if isinstance(item, dict):
                caption = _coerce_text(item.get('caption') or item.get('alt_text') or item.get('title') or item.get('slot_key'))
                image_path = _coerce_text(
                    item.get('image_path')
                    or item.get('path')
                    or item.get('file_path')
                    or item.get('local_path')
                )
                images.append(
                    {
                        'slot_key': _coerce_text(item.get('slot_key') or item.get('key') or caption),
                        'caption': caption,
                        'alt_text': _coerce_text(item.get('alt_text') or caption),
                        'source': _coerce_text(item.get('source') or item.get('content_source')),
                        'image_path': image_path,
                        'path': image_path,
                    }
                )
            else:
                caption = _coerce_text(item)
                if caption:
                    images.append(
                        {
                            'slot_key': caption,
                            'caption': caption,
                            'alt_text': caption,
                            'source': '',
                            'image_path': '',
                            'path': '',
                        }
                    )

    return images


def _add_title_box(slide, title: str, *, page_function: str, style: dict[str, Any]):
    title_y = Inches(0.35 if page_function != 'cover' else 1.1)
    title_h = Inches(0.8 if page_function != 'cover' else 1.2)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.8))
    title_box.left = Inches(0.6)
    title_box.top = title_y
    title_box.width = Inches(12.0)
    title_box.height = title_h
    text_frame = title_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    font = run.font
    font.name = str(style.get('font_family') or 'Arial')
    font.size = Pt(34 if page_function == 'cover' else 28)
    font.bold = True
    font.color.rgb = _safe_hex_color(style.get('text_color'), (30, 43, 57))
    paragraph.alignment = PP_ALIGN.CENTER if page_function in {'cover', 'ending'} else PP_ALIGN.LEFT
    return title_box


def _add_bullets_box(
    slide,
    bullets: list[str],
    *,
    page_function: str,
    style: dict[str, Any],
    asset_mode: bool = False,
):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.4), Inches(5.9))
    if page_function == 'cover':
        box.top = Inches(2.5)
        box.height = Inches(3.8)
    elif page_function == 'toc':
        box.top = Inches(1.55)
        box.height = Inches(3.6 if asset_mode else 5.5)
    elif page_function == 'comparison':
        box.width = Inches(5.4 if not asset_mode else 5.0)
        box.height = Inches(5.6 if not asset_mode else 3.1)
    elif asset_mode:
        box.height = Inches(2.95)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Pt(2)
    text_frame.margin_right = Pt(2)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.level = 0
        run = paragraph.add_run()
        prefix = f'{idx + 1}. ' if page_function == 'toc' else ''
        run.text = f'{prefix}{bullet}'
        run.font.name = str(style.get('font_family') or 'Arial')
        run.font.size = Pt(16 if page_function == 'toc' else 18)
        run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(4)

    if page_function == 'comparison' and not asset_mode:
        right_box = slide.shapes.add_textbox(Inches(6.95), Inches(1.35), Inches(5.4), Inches(5.6))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        right_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        split = max(1, len(bullets) // 2)
        right_items = bullets[split:]
        if not right_items:
            right_items = ['Comparison note', 'Trade-offs', 'Recommendation']
        for idx, bullet in enumerate(right_items):
            paragraph = right_tf.paragraphs[0] if idx == 0 else right_tf.add_paragraph()
            paragraph.level = 0
            run = paragraph.add_run()
            run.text = bullet
            run.font.name = str(style.get('font_family') or 'Arial')
            run.font.size = Pt(18)
            run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(4)
    return box


def _apply_background(slide, style: dict[str, Any]):
    bg = _safe_hex_color(style.get('background_color'), (247, 250, 255))
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg


def _add_accent_bar(slide, style: dict[str, Any]):
    accent = _safe_hex_color(style.get('accent_color'), (27, 110, 243))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(2.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def _add_table_placeholder(
    slide,
    table_data: dict[str, Any],
    *,
    left,
    top,
    width,
    style: dict[str, Any],
):
    headers = [_coerce_text(value, default='') for value in _coerce_items(table_data.get('headers'))]
    rows_raw = _coerce_items(table_data.get('rows'))
    rows: list[list[str]] = []
    for row in rows_raw:
        if isinstance(row, dict):
            rows.append([_coerce_text(value, default='-') for value in row.values()])
        elif isinstance(row, (list, tuple)):
            rows.append([_coerce_text(value, default='-') for value in row])
        else:
            row_text = _coerce_text(row, default='-')
            if row_text:
                rows.append([row_text])

    if not rows:
        rows = [['No data']]

    column_count = max(1, len(headers), max(len(row) for row in rows))
    normalized_headers = headers[:column_count] + [''] * max(0, column_count - len(headers))
    normalized_rows: list[list[str]] = []
    for row in rows[:6]:
        normalized_rows.append(row[:column_count] + [''] * max(0, column_count - len(row)))

    row_count = len(normalized_rows) + 1
    table_height = Inches(min(2.5, max(1.15, 0.28 * row_count + 0.38)))
    table_shape = slide.shapes.add_table(row_count, column_count, left, top, width, table_height)
    table = table_shape.table
    column_width = int(width / column_count)
    accent = _safe_hex_color(style.get('accent_color'), (27, 110, 243))
    text_color = _safe_hex_color(style.get('text_color'), (35, 49, 66))

    for col_idx in range(column_count):
        table.columns[col_idx].width = column_width

    for col_idx, header in enumerate(normalized_headers):
        cell = table.cell(0, col_idx)
        cell.text = header or f'Column {col_idx + 1}'
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = str(style.get('font_family') or 'Arial')
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, row in enumerate(normalized_rows, start=1):
        for col_idx in range(column_count):
            cell = table.cell(row_idx, col_idx)
            cell.text = row[col_idx] if col_idx < len(row) else ''
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = str(style.get('font_family') or 'Arial')
                    run.font.size = Pt(11)
                    run.font.color.rgb = text_color

    return table_shape


def _table_text_length(table_data: dict[str, Any]) -> int:
    total = 0
    for value in _coerce_items(table_data.get('headers')):
        total += len(_coerce_text(value))
    for row in _coerce_items(table_data.get('rows')):
        if isinstance(row, dict):
            total += sum(len(_coerce_text(value)) for value in row.values())
        elif isinstance(row, (list, tuple)):
            total += sum(len(_coerce_text(value)) for value in row)
        else:
            total += len(_coerce_text(row))
    return total


def _is_complex_table(table_data: dict[str, Any]) -> bool:
    headers = [_coerce_text(value, default='') for value in _coerce_items(table_data.get('headers')) if _coerce_text(value)]
    rows = _coerce_items(table_data.get('rows'))
    if len(headers) >= 5:
        return True
    if len(rows) >= 6:
        return True
    if _table_text_length(table_data) >= 240:
        return True
    for row in rows:
        if isinstance(row, dict):
            values = list(row.values())
        elif isinstance(row, (list, tuple)):
            values = list(row)
        else:
            values = [row]
        if any(len(_coerce_text(value)) >= 36 for value in values):
            return True
    return False


def _add_table_summary_card(
    slide,
    card_data: dict[str, Any],
    *,
    left,
    top,
    width,
    style: dict[str, Any],
):
    accent = _safe_hex_color(style.get('accent_color'), (27, 110, 243))
    text_color = _safe_hex_color(style.get('text_color'), (35, 49, 66))
    background = RGBColor(245, 247, 250)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.72))
    box.fill.solid()
    box.fill.fore_color.rgb = background
    box.line.color.rgb = accent
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Pt(5)
    box.text_frame.margin_right = Pt(5)
    box.text_frame.margin_top = Pt(4)
    box.text_frame.margin_bottom = Pt(4)

    title = _coerce_text(card_data.get('title') or 'Table Summary')
    headers = [_coerce_text(value, default='') for value in _coerce_items(card_data.get('headers')) if _coerce_text(value)]
    rows = _coerce_items(card_data.get('rows'))
    row_count = len(rows)
    preview_items: list[str] = []
    for row in rows[:3]:
        if isinstance(row, dict):
            preview_items.append(' | '.join(_coerce_text(value) for value in row.values() if _coerce_text(value)))
        elif isinstance(row, (list, tuple)):
            preview_items.append(' | '.join(_coerce_text(value) for value in row if _coerce_text(value)))
        else:
            preview_items.append(_coerce_text(row))
    preview = '; '.join([item for item in preview_items if item])[:160]
    if not preview:
        preview = 'Summary of complex table content.'
    header_summary = ', '.join(headers[:4]) if headers else 'table data'

    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = f'{title}'
    run.font.name = str(style.get('font_family') or 'Arial')
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = accent
    paragraph.alignment = PP_ALIGN.LEFT

    body = box.text_frame.add_paragraph()
    body.text = f'{header_summary}. {preview}'
    body.alignment = PP_ALIGN.LEFT
    for run in body.runs:
        run.font.name = str(style.get('font_family') or 'Arial')
        run.font.size = Pt(10)
        run.font.color.rgb = text_color

    meta = box.text_frame.add_paragraph()
    meta.text = f'Complex table downgraded: {row_count} rows'
    meta.alignment = PP_ALIGN.LEFT
    for run in meta.runs:
        run.font.name = str(style.get('font_family') or 'Arial')
        run.font.size = Pt(9)
        run.font.color.rgb = text_color
    return box


def _resolve_local_image_path(image_data: dict[str, Any]) -> Path | None:
    candidate = _coerce_text(
        image_data.get('image_path')
        or image_data.get('path')
        or image_data.get('file_path')
        or image_data.get('local_path')
    )
    if not candidate:
        return None

    try:
        image_path = Path(candidate).expanduser()
    except Exception:
        return None

    allowed_suffixes = {'.png', '.jpg', '.jpeg', '.bmp', '.gif'}
    if image_path.suffix.lower() not in allowed_suffixes:
        return None
    if not image_path.is_file():
        return None
    return image_path


def _fit_picture_box(
    image_path: Path | None,
    *,
    left,
    top,
    width,
    max_height,
) -> tuple[int, int, int, int] | None:
    if image_path is None:
        return None

    picture_width = int(width)
    picture_height = None
    picture_left = int(left)
    picture_top = int(top)
    try:
        from PIL import Image as PILImage

        with PILImage.open(image_path) as image:
            pixel_width, pixel_height = image.size
        if pixel_width > 0 and pixel_height > 0:
            max_width_emu = int(width)
            max_height_emu = int(max_height)
            aspect_ratio = pixel_width / pixel_height
            box_ratio = max_width_emu / max_height_emu if max_height_emu else aspect_ratio
            if aspect_ratio >= box_ratio:
                picture_width = max_width_emu
                picture_height = int(max_width_emu / aspect_ratio)
            else:
                picture_height = max_height_emu
                picture_width = int(max_height_emu * aspect_ratio)
            picture_left = int(left + max(0, int((max_width_emu - picture_width) / 2)))
            picture_top = int(top + max(0, int((max_height_emu - picture_height) / 2)))
    except Exception:
        picture_width = int(width)
        picture_height = None
        picture_left = int(left)
        picture_top = int(top)

    if picture_height is None:
        return picture_left, picture_top, picture_width, 0
    return picture_left, picture_top, picture_width, picture_height


def _try_add_picture(
    slide,
    image_path: Path | None,
    *,
    left,
    top,
    width,
    max_height,
):
    if image_path is None:
        return None

    fitted = _fit_picture_box(image_path, left=left, top=top, width=width, max_height=max_height)
    if fitted is None:
        return None
    picture_left, picture_top, picture_width, picture_height = fitted

    try:
        if picture_height <= 0:
            return slide.shapes.add_picture(str(image_path), picture_left, picture_top, width=picture_width)
        return slide.shapes.add_picture(str(image_path), picture_left, picture_top, width=picture_width, height=picture_height)
    except Exception:
        return None


def _add_image_asset(
    slide,
    image_data: dict[str, Any],
    *,
    left,
    top,
    width,
    max_height,
    style: dict[str, Any],
):
    accent = _safe_hex_color(style.get('accent_color'), (27, 110, 243))
    caption_text = _coerce_text(image_data.get('caption') or image_data.get('alt_text') or 'Image')
    image_path = _resolve_local_image_path(image_data)
    picture = _try_add_picture(slide, image_path, left=left, top=top, width=width, max_height=max_height)

    if picture is not None:
        caption_top = top + picture.height + Inches(0.08)
        caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.36))
        caption_tf = caption_box.text_frame
        caption_tf.word_wrap = True
        caption_tf.clear()
        paragraph = caption_tf.paragraphs[0]
        run = paragraph.add_run()
        run.text = caption_text
        run.font.name = str(style.get('font_family') or 'Arial')
        run.font.size = Pt(11)
        run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
        paragraph.alignment = PP_ALIGN.CENTER
        return caption_box

    frame_height = Inches(1.35)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, frame_height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(245, 247, 250)
    frame.line.color.rgb = accent
    frame.text_frame.clear()
    frame.text_frame.word_wrap = True
    frame.text_frame.margin_left = Pt(4)
    frame.text_frame.margin_right = Pt(4)
    frame.text_frame.margin_top = Pt(4)
    frame.text_frame.margin_bottom = Pt(4)
    paragraph = frame.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = 'IMAGE PLACEHOLDER'
    run.font.name = str(style.get('font_family') or 'Arial')
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = accent
    paragraph.alignment = PP_ALIGN.CENTER

    caption_box = slide.shapes.add_textbox(left, top + frame_height + Inches(0.03), width, Inches(0.36))
    caption_tf = caption_box.text_frame
    caption_tf.word_wrap = True
    caption_tf.clear()
    paragraph = caption_tf.paragraphs[0]
    run = paragraph.add_run()
    run.text = caption_text
    run.font.name = str(style.get('font_family') or 'Arial')
    run.font.size = Pt(11)
    run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
    paragraph.alignment = PP_ALIGN.CENTER

    return caption_box


def _render_assets(
    slide,
    *,
    tables: list[dict[str, Any]],
    images: list[dict[str, Any]],
    downgraded_tables: list[dict[str, Any]] | None = None,
    page_function: str,
    style: dict[str, Any],
):
    has_tables = bool(tables)
    has_images = bool(images)
    has_downgraded_tables = bool(downgraded_tables)
    if not (has_tables or has_images or has_downgraded_tables):
        return {'fallback_tables': 0, 'fallback_images': 0}

    zone_top = Inches(4.15 if page_function != 'cover' else 4.65)
    gap = Inches(0.18)
    full_left = Inches(0.75)
    full_width = Inches(11.85)
    half_width = Inches(5.72)
    left_column = full_left
    right_column = Inches(6.82)

    if has_tables and has_images:
        table_left = left_column
        table_width = half_width
        image_left = right_column
        image_width = half_width
    else:
        table_left = full_left
        table_width = full_width
        image_left = full_left
        image_width = full_width

    current_top = zone_top
    fallback_tables = 0
    if has_tables:
        for table_data in tables[:3]:
            rendered = _add_table_placeholder(
                slide,
                table_data,
                left=table_left,
                top=current_top,
                width=table_width,
                style=style,
            )
            current_top = rendered.top + rendered.height + gap
            fallback_tables += 1

    if has_downgraded_tables:
        for card_data in downgraded_tables[:3]:
            rendered = _add_table_summary_card(
                slide,
                card_data,
                left=table_left,
                top=current_top,
                width=table_width,
                style=style,
            )
            current_top = rendered.top + rendered.height + gap
            fallback_tables += 1

    current_top = zone_top
    fallback_images = 0
    if has_images:
        for image_data in images[:3]:
            rendered = _add_image_asset(
                slide,
                image_data,
                left=image_left,
                top=current_top,
                width=image_width,
                max_height=Inches(1.95 if has_tables else 2.25),
                style=style,
            )
            current_top = rendered.top + rendered.height + gap
            fallback_images += 1
    return {'fallback_tables': fallback_tables, 'fallback_images': fallback_images}


def _clear_presentation_slides(prs: Presentation) -> None:
    slide_id_list = list(prs.slides._sldIdLst)
    for slide_id in slide_id_list:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def _resolve_template_pair(template_path: Path | None) -> tuple[Presentation | None, Presentation]:
    if template_path and template_path.is_file() and template_path.suffix.lower() == '.pptx':
        source_prs = Presentation(str(template_path))
        output_prs = Presentation(str(template_path))
        _clear_presentation_slides(output_prs)
        return source_prs, output_prs

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return None, prs


def _layout_index(source_prs: Presentation, layout_obj: Any) -> int:
    for idx, candidate in enumerate(source_prs.slide_layouts):
        if candidate == layout_obj:
            return idx
    return 0


def _clone_template_shapes(source_slide, target_slide) -> int:
    sp_tree = target_slide.shapes._spTree
    for shape in list(target_slide.shapes):
        try:
            sp_tree.remove(shape.element)
        except Exception:
            continue

    skipped_types = {
        getattr(MSO_SHAPE_TYPE, 'PICTURE', None),
        getattr(MSO_SHAPE_TYPE, 'CHART', None),
        getattr(MSO_SHAPE_TYPE, 'MEDIA', None),
        getattr(MSO_SHAPE_TYPE, 'EMBEDDED_OLE_OBJECT', None),
        getattr(MSO_SHAPE_TYPE, 'LINKED_OLE_OBJECT', None),
        getattr(MSO_SHAPE_TYPE, 'LINKED_PICTURE', None),
    }
    copied = 0
    for shape in source_slide.shapes:
        try:
            if shape.shape_type in skipped_types:
                continue
            sp_tree.insert_element_before(deepcopy(shape.element), 'p:extLst')
            copied += 1
        except Exception:
            continue
    return copied


def _add_slide_from_template(
    prs: Presentation,
    template_source_prs: Presentation | None,
    slide_data: dict[str, Any],
    slide_index: int,
):
    if template_source_prs is None:
        blank_layout_idx = min(6, len(prs.slide_layouts) - 1)
        return prs.slides.add_slide(prs.slide_layouts[blank_layout_idx]), False

    template_page_no = int(slide_data.get('template_page_no') or slide_data.get('page_no') or slide_index)
    if template_page_no < 1 or template_page_no > len(template_source_prs.slides):
        blank_layout_idx = min(6, len(prs.slide_layouts) - 1)
        return prs.slides.add_slide(prs.slide_layouts[blank_layout_idx]), False

    source_slide = template_source_prs.slides[template_page_no - 1]
    layout_idx = _layout_index(template_source_prs, source_slide.slide_layout)
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
    new_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    _clone_template_shapes(source_slide, new_slide)
    return new_slide, True


def _set_shape_title(shape, title: str, style: dict[str, Any]) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = str(style.get('font_family') or 'Arial')
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = _safe_hex_color(style.get('text_color'), (30, 43, 57))
    paragraph.alignment = PP_ALIGN.LEFT


def _set_shape_bullets(shape, bullets: list[str], style: dict[str, Any]) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = bullet
        run.font.name = str(style.get('font_family') or 'Arial')
        run.font.size = Pt(18)
        run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(4)


def _set_shape_subtitle(shape, subtitle: str, style: dict[str, Any]) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = subtitle
    run.font.name = str(style.get('font_family') or 'Arial')
    run.font.size = Pt(18)
    run.font.bold = False
    run.font.color.rgb = _safe_hex_color(style.get('text_color'), (35, 49, 66))
    paragraph.alignment = PP_ALIGN.LEFT


def _shape_position_key(shape) -> tuple[int, int]:
    return (int(getattr(shape, 'top', 0)), int(getattr(shape, 'left', 0)))


def _slot_key_tokens(slot_key: str) -> list[str]:
    key = _coerce_text(slot_key).lower()
    if not key:
        return []
    return [token for token in re.split(r"[^a-z0-9]+", key) if len(token) >= 2]


def _normalize_layout_slots_for_template(slide_data: dict[str, Any]) -> list[dict[str, str]]:
    layout_schema = slide_data.get('layout_schema_json')
    if not isinstance(layout_schema, dict):
        return []

    raw_slots = layout_schema.get('slots') or []
    normalized: list[dict[str, str]] = []
    for idx, item in enumerate(_coerce_items(raw_slots), start=1):
        if isinstance(item, dict):
            slot_key = _coerce_text(item.get('slot_key') or item.get('key'), default=f'slot_{idx}')
            slot_type = _coerce_text(item.get('slot_type'), default='text').lower()
            slot_role = _coerce_text(item.get('slot_role') or item.get('role')).lower()
        else:
            slot_key = _coerce_text(item, default=f'slot_{idx}')
            lowered = slot_key.lower()
            slot_type = 'table' if 'table' in lowered else ('image' if any(k in lowered for k in ('image', 'figure', 'photo', 'pic')) else 'text')
            if 'title' in lowered:
                slot_role = 'title'
            elif 'subtitle' in lowered:
                slot_role = 'subtitle'
            elif slot_type == 'table':
                slot_role = 'datatable'
            elif slot_type == 'image':
                slot_role = 'figure'
            else:
                slot_role = 'bullet'
        normalized.append(
            {
                'slot_key': slot_key,
                'slot_type': slot_type,
                'slot_role': slot_role or 'bullet',
            }
        )
    return normalized


def _text_shape_score(shape, *, slot_role: str, slot_key: str) -> int:
    score = 0
    try:
        placeholder_type = shape.placeholder_format.type if shape.is_placeholder else None
    except Exception:
        placeholder_type = None
    shape_name = str(getattr(shape, 'name', '')).lower()
    name_tokens = _slot_key_tokens(shape_name)
    slot_tokens = _slot_key_tokens(slot_key)
    for token in slot_tokens:
        if token in name_tokens or token in shape_name:
            score += 12

    if slot_role == 'title':
        if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
            score += 120
        if 'title' in shape_name:
            score += 36
    elif slot_role == 'subtitle':
        if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
            score += 120
        if 'subtitle' in shape_name:
            score += 30
    else:
        if placeholder_type in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}:
            score += 80
        if any(keyword in shape_name for keyword in ('content', 'body', 'text', 'bullet')):
            score += 25
    return score


def _pick_best_text_shape(
    text_shapes: list[Any],
    *,
    slot_role: str,
    slot_key: str,
    used_ids: set[int],
) -> Any | None:
    ranked: list[tuple[int, int, Any]] = []
    for index, shape in enumerate(text_shapes):
        shape_id = id(shape)
        if shape_id in used_ids:
            continue
        score = _text_shape_score(shape, slot_role=slot_role, slot_key=slot_key)
        ranked.append((score, -index, shape))
    if not ranked:
        return None
    ranked.sort(reverse=True)
    return ranked[0][2]


def _apply_template_text_edits(
    slide,
    title: str,
    bullets: list[str],
    style: dict[str, Any],
    slide_data: dict[str, Any],
) -> dict[str, Any]:
    title_done = False
    body_done = False
    subtitle_done = False
    subtitle_text = bullets[0] if bullets else title
    text_ops = 0
    slot_named_matches = 0
    text_shapes = [shape for shape in sorted(slide.shapes, key=_shape_position_key) if getattr(shape, 'has_text_frame', False)]
    used_text_shapes: set[int] = set()

    slot_specs = [slot for slot in _normalize_layout_slots_for_template(slide_data) if slot.get('slot_type') == 'text']
    if not slot_specs:
        slot_specs = [
            {'slot_key': 'title', 'slot_type': 'text', 'slot_role': 'title'},
            {'slot_key': 'subtitle', 'slot_type': 'text', 'slot_role': 'subtitle'},
            {'slot_key': 'body', 'slot_type': 'text', 'slot_role': 'bullet'},
        ]

    for slot in slot_specs:
        role = str(slot.get('slot_role') or 'bullet').lower()
        key = str(slot.get('slot_key') or role)
        candidate = _pick_best_text_shape(text_shapes, slot_role=role, slot_key=key, used_ids=used_text_shapes)
        if candidate is None:
            continue
        used_text_shapes.add(id(candidate))
        candidate_name = str(getattr(candidate, 'name', '')).lower()
        if any(token and token in candidate_name for token in _slot_key_tokens(key)):
            slot_named_matches += 1

        if role == 'title' and not title_done:
            _set_shape_title(candidate, title, style)
            title_done = True
            text_ops += 1
            continue
        if role == 'subtitle' and not subtitle_done and subtitle_text:
            _set_shape_subtitle(candidate, subtitle_text, style)
            subtitle_done = True
            text_ops += 1
            continue
        if not body_done:
            _set_shape_bullets(candidate, bullets, style)
            body_done = True
            text_ops += 1

    if not title_done:
        for shape in text_shapes:
            if id(shape) in used_text_shapes:
                continue
            current_text = _coerce_text(getattr(shape, 'text', ''))
            if 'title' in _coerce_text(getattr(shape, 'name', '')).lower() or (current_text and len(current_text) <= 80):
                _set_shape_title(shape, title, style)
                title_done = True
                text_ops += 1
                used_text_shapes.add(id(shape))
                break

    if not body_done:
        for shape in text_shapes:
            if id(shape) in used_text_shapes:
                continue
            current_text = _coerce_text(getattr(shape, 'text', ''))
            if title and current_text == title:
                continue
            _set_shape_bullets(shape, bullets, style)
            body_done = True
            text_ops += 1
            break
    return {
        'title_done': title_done,
        'subtitle_done': subtitle_done,
        'body_done': body_done,
        'text_ops': text_ops,
        'slot_named_matches': slot_named_matches,
    }


def _apply_table_data_to_shape(shape, table_data: dict[str, Any], style: dict[str, Any]) -> None:
    table = shape.table
    total_rows = len(table.rows)
    total_cols = len(table.columns)
    if total_rows <= 0 or total_cols <= 0:
        return

    headers = [_coerce_text(value, default='') for value in _coerce_items(table_data.get('headers'))]
    rows = []
    for row in _coerce_items(table_data.get('rows')):
        if isinstance(row, (list, tuple)):
            rows.append([_coerce_text(value, default='') for value in row])
        elif isinstance(row, dict):
            rows.append([_coerce_text(value, default='') for value in row.values()])
        else:
            rows.append([_coerce_text(row, default='')])

    headers = (headers + [''] * total_cols)[:total_cols]
    rows = [((row + [''] * total_cols)[:total_cols]) for row in rows[: max(0, total_rows - 1)]]

    accent = _safe_hex_color(style.get('accent_color'), (27, 110, 243))
    text_color = _safe_hex_color(style.get('text_color'), (35, 49, 66))
    for col_idx in range(total_cols):
        cell = table.cell(0, col_idx)
        cell.text = headers[col_idx] or f'Column {col_idx + 1}'
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = str(style.get('font_family') or 'Arial')
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx in range(1, total_rows):
        row_values = rows[row_idx - 1] if row_idx - 1 < len(rows) else [''] * total_cols
        for col_idx in range(total_cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = row_values[col_idx]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = str(style.get('font_family') or 'Arial')
                    run.font.size = Pt(10)
                    run.font.color.rgb = text_color


def _pick_best_table_shape(table_shapes: list[Any], slot_key: str, used_ids: set[int]) -> Any | None:
    best_shape = None
    best_score = -1
    tokens = _slot_key_tokens(slot_key)
    for index, shape in enumerate(table_shapes):
        shape_id = id(shape)
        if shape_id in used_ids:
            continue
        shape_name = _coerce_text(getattr(shape, 'name', '')).lower()
        score = 50
        if 'table' in shape_name:
            score += 20
        for token in tokens:
            if token and token in shape_name:
                score += 12
        score += max(0, 8 - index)
        if score > best_score:
            best_score = score
            best_shape = shape
    return best_shape


def _apply_template_table_edits(slide, tables: list[dict[str, Any]], style: dict[str, Any]) -> dict[str, int]:
    consumed = 0
    named_matches = 0
    if not tables:
        return {'consumed': 0, 'named_matches': 0}

    table_shapes = [shape for shape in sorted(slide.shapes, key=_shape_position_key) if getattr(shape, 'has_table', False)]
    used_shapes: set[int] = set()
    for table in tables:
        slot_key = _coerce_text(table.get('slot_key') or table.get('title') or 'table')
        target = _pick_best_table_shape(table_shapes, slot_key, used_shapes)
        if target is None:
            continue
        used_shapes.add(id(target))
        _apply_table_data_to_shape(target, table, style)
        consumed += 1
        shape_name = _coerce_text(getattr(target, 'name', '')).lower()
        if any(token and token in shape_name for token in _slot_key_tokens(slot_key)):
            named_matches += 1
    return {'consumed': consumed, 'named_matches': named_matches}


def _pick_best_image_shape(image_shapes: list[Any], slot_key: str, used_ids: set[int]) -> Any | None:
    best_shape = None
    best_score = -1
    tokens = _slot_key_tokens(slot_key)
    for index, shape in enumerate(image_shapes):
        shape_id = id(shape)
        if shape_id in used_ids:
            continue
        shape_name = _coerce_text(getattr(shape, 'name', '')).lower()
        score = 40
        try:
            placeholder_type = shape.placeholder_format.type if shape.is_placeholder else None
        except Exception:
            placeholder_type = None
        if placeholder_type == PP_PLACEHOLDER.PICTURE:
            score += 40
        if any(keyword in shape_name for keyword in ('image', 'picture', 'photo', 'figure', 'pic')):
            score += 24
        for token in tokens:
            if token and token in shape_name:
                score += 10
        score += max(0, 8 - index)
        if score > best_score:
            best_score = score
            best_shape = shape
    return best_shape


def _apply_template_image_edits(slide, images: list[dict[str, Any]]) -> dict[str, int]:
    consumed = 0
    named_matches = 0
    overlay_inserts = 0
    if not images:
        return {'consumed': 0, 'named_matches': 0, 'overlay_inserts': 0}

    image_shapes: list[Any] = []
    for shape in sorted(slide.shapes, key=_shape_position_key):
        try:
            placeholder_type = shape.placeholder_format.type if shape.is_placeholder else None
        except Exception:
            placeholder_type = None
        shape_name = _coerce_text(getattr(shape, 'name', '')).lower()
        if placeholder_type == PP_PLACEHOLDER.PICTURE or any(keyword in shape_name for keyword in ('image', 'picture', 'photo', 'figure', 'pic')):
            image_shapes.append(shape)

    used_shapes: set[int] = set()
    for image_data in images:
        image_path = _resolve_local_image_path(image_data)
        if image_path is None:
            continue
        slot_key = _coerce_text(image_data.get('slot_key') or image_data.get('caption') or image_data.get('alt_text') or 'image')
        target = _pick_best_image_shape(image_shapes, slot_key, used_shapes)
        if target is None:
            continue

        inserted = False
        fitted_box = _fit_picture_box(
            image_path,
            left=target.left,
            top=target.top,
            width=target.width,
            max_height=target.height,
        )
        insert_picture = getattr(target, 'insert_picture', None)
        if callable(insert_picture):
            try:
                picture = insert_picture(str(image_path))
                if picture is not None and fitted_box is not None:
                    fitted_left, fitted_top, fitted_width, fitted_height = fitted_box
                    picture.left = fitted_left
                    picture.top = fitted_top
                    picture.width = fitted_width
                    if fitted_height > 0:
                        picture.height = fitted_height
                inserted = True
            except Exception:
                inserted = False
        if not inserted:
            picture = _try_add_picture(
                slide,
                image_path,
                left=target.left,
                top=target.top,
                width=target.width,
                max_height=target.height,
            )
            if picture is not None:
                inserted = True
                overlay_inserts += 1
        if not inserted:
            continue

        used_shapes.add(id(target))
        consumed += 1
        shape_name = _coerce_text(getattr(target, 'name', '')).lower()
        if any(token and token in shape_name for token in _slot_key_tokens(slot_key)):
            named_matches += 1

    return {'consumed': consumed, 'named_matches': named_matches, 'overlay_inserts': overlay_inserts}


def generate_pptx_from_plan(
    slide_plan: list[dict],
    output_path: Path,
    title_prefix: str = '',
    template_path: Path | None = None,
) -> dict:
    """
    Render a slide plan into an editable PPTX deck.

    Each slide gets a title and a bullet list, with deterministic fallback
    content when the input omits bullets.
    """

    try:
        if not isinstance(output_path, Path):
            output_path = Path(output_path)

        if not slide_plan:
            raise PPTXGenerationError('slide_plan is empty')

        output_path.parent.mkdir(parents=True, exist_ok=True)

        template_source_prs, prs = _resolve_template_pair(template_path)

        rendered_titles: list[str] = []
        template_text_ops = 0
        template_table_ops = 0
        template_image_ops = 0
        template_named_matches = 0
        template_overlay_image_inserts = 0
        slide_edit_stats: list[dict[str, Any]] = []
        for index, slide_data in enumerate(slide_plan, start=1):
            if not isinstance(slide_data, dict):
                raise PPTXGenerationError(f'slide_plan[{index - 1}] must be a dict')

            slide, template_applied = _add_slide_from_template(prs, template_source_prs, slide_data, index)
            page_function = _coerce_text(slide_data.get('page_function'), default='content').lower()
            style = _style_tokens(slide_data)
            title = _normalize_title(slide_data, index, title_prefix)
            bullets = _normalize_bullets(slide_data, index)
            tables = _normalize_tables(slide_data)
            render_tables: list[dict[str, Any]] = []
            downgraded_tables: list[dict[str, Any]] = []
            for table in tables:
                if _is_complex_table(table):
                    downgraded_tables.append(table)
                else:
                    render_tables.append(table)
            images = _normalize_images(slide_data)
            asset_mode = bool(render_tables or downgraded_tables or images)
            slide_stat = {
                'slide_no': int(slide_data.get('page_no') or index),
                'template_page_no': int(slide_data.get('template_page_no') or slide_data.get('page_no') or index),
                'template_applied': bool(template_applied),
                'text_ops': 0,
                'table_ops': 0,
                'image_ops': 0,
                'named_matches': 0,
                'fallback_tables': 0,
                'fallback_images': 0,
                'image_overlay_inserts': 0,
                'downgraded_table_count': len(downgraded_tables),
                'image_fit_strategy': 'contain_center' if images else 'none',
            }

            if template_applied:
                text_info = _apply_template_text_edits(slide, title, bullets, style, slide_data)
                table_info = _apply_template_table_edits(slide, render_tables, style)
                image_info = _apply_template_image_edits(slide, images)
                template_text_ops += int(text_info.get('text_ops') or 0)
                template_table_ops += int(table_info.get('consumed') or 0)
                template_image_ops += int(image_info.get('consumed') or 0)
                template_named_matches += (
                    int(text_info.get('slot_named_matches') or 0)
                    + int(table_info.get('named_matches') or 0)
                    + int(image_info.get('named_matches') or 0)
                )
                template_overlay_image_inserts += int(image_info.get('overlay_inserts') or 0)

                if not text_info.get('title_done'):
                    _add_title_box(slide, title, page_function=page_function, style=style)
                if not text_info.get('body_done'):
                    _add_bullets_box(slide, bullets, page_function=page_function, style=style, asset_mode=asset_mode)
                remaining_tables = render_tables[int(table_info.get('consumed') or 0) :]
                remaining_downgraded_tables = downgraded_tables
                remaining_images = images[int(image_info.get('consumed') or 0) :]
                slide_stat['text_ops'] = int(text_info.get('text_ops') or 0)
                slide_stat['table_ops'] = int(table_info.get('consumed') or 0)
                slide_stat['image_ops'] = int(image_info.get('consumed') or 0)
                slide_stat['named_matches'] = (
                    int(text_info.get('slot_named_matches') or 0)
                    + int(table_info.get('named_matches') or 0)
                    + int(image_info.get('named_matches') or 0)
                )
                slide_stat['image_overlay_inserts'] = int(image_info.get('overlay_inserts') or 0)
            else:
                _apply_background(slide, style)
                _add_title_box(slide, title, page_function=page_function, style=style)
                _add_accent_bar(slide, style)
                _add_bullets_box(slide, bullets, page_function=page_function, style=style, asset_mode=asset_mode)
                remaining_tables = render_tables
                remaining_downgraded_tables = downgraded_tables
                remaining_images = images
            fallback_info = _render_assets(
                slide,
                tables=remaining_tables,
                images=remaining_images,
                downgraded_tables=remaining_downgraded_tables,
                page_function=page_function,
                style=style,
            )
            slide_stat['fallback_tables'] = int(fallback_info.get('fallback_tables') or 0)
            slide_stat['fallback_images'] = int(fallback_info.get('fallback_images') or 0)
            slide_edit_stats.append(slide_stat)
            rendered_titles.append(title)

        prs.save(output_path)
        file_size = output_path.stat().st_size

        return {
            'output_path': str(output_path),
            'page_count': len(slide_plan),
            'file_size': file_size,
            'titles': rendered_titles,
            'template_mode': bool(template_source_prs),
            'template_edit_stats': {
                'text_placeholder_edits': template_text_ops,
                'table_placeholder_edits': template_table_ops,
                'image_placeholder_edits': template_image_ops,
                'named_slot_matches': template_named_matches,
                'image_overlay_inserts': template_overlay_image_inserts,
                'slide_stats': slide_edit_stats,
            },
        }
    except PPTXGenerationError:
        raise
    except Exception as exc:  # pragma: no cover - defensive wrapper
        raise PPTXGenerationError(f'failed to generate pptx: {exc}') from exc

